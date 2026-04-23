#!/usr/bin/env python3
"""Build docs/PartSelect-Design-Deck.pptx from DESIGN.md themes (run: python3 docs/generate_design_deck.py)."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs, title: str, subtitle: str = ""):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_bullets(prs, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def add_two_col_table(prs, title: str, rows: list[tuple[str, str]]):
    layout = prs.slide_layouts[5]  # blank
    slide = prs.slides.add_slide(layout)
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(28)
    tx.text_frame.paragraphs[0].font.bold = True

    rows_n = len(rows) + 1
    cols = 2
    left, top, width, height = Inches(0.5), Inches(1.1), Inches(9), Inches(4.8)
    table = slide.shapes.add_table(rows_n, cols, left, top, width, height).table
    table.columns[0].width = Inches(4.2)
    table.columns[1].width = Inches(4.8)
    hdr = table.rows[0].cells
    hdr[0].text = "In scope"
    hdr[1].text = "Out of scope"
    for c in hdr:
        c.text_frame.paragraphs[0].font.bold = True
    for i, (a, b) in enumerate(rows, start=1):
        table.rows[i].cells[0].text = a
        table.rows[i].cells[1].text = b


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "PartSelect Chat Agent",
        "Architecture & design (from DESIGN.md)\nRefrigerator & dishwasher parts",
    )

    add_bullets(
        prs,
        "Summary",
        [
            "Domain-scoped assistant: fridge & dishwasher parts only",
            "POST /api/chat — NDJSON stream (token / replace → done), not SSE",
            "retrieveExact + catalog.json; optional semantic_search (embeddings)",
            "UI: chat (cards + chips) + left mock cart (local demo → PartSelect checkout)",
        ],
    )

    add_bullets(
        prs,
        "Goals",
        [
            "In-domain: lookup, price/stock, OEM/supersession, install, compat, symptom → repair/candidates",
            "Refuse out-of-domain clearly",
            "Ground copy & cards in catalog + citations",
            "Stream replies; optional RAG for experiential questions",
        ],
    )

    add_two_col_table(
        prs,
        "Scope",
        [
            ("PS / OEM / symptoms / install / compat", "Orders, returns, live handoff"),
            ("catalog.json + optional PDP scrape", "Washers, HVAC… (refused)"),
            ("Read-only PartSelect HTML (fetch_part_page, images)", "Official PartSelect API"),
        ],
    )

    add_bullets(
        prs,
        "Cross-cutting (2)",
        [
            "Server-built blocks — buildBlocksFromRetrieval; no fragile LLM JSON for card payloads",
            "Stable done JSON: reply, blocks, citations, suggested_actions, tool_trace, no_evidence",
        ],
    )

    add_bullets(
        prs,
        "Interface",
        [
            "NDJSON over POST — message + history in body (EventSource / SSE is GET-only)",
            "Cards + chips — structured UI; clarify tips stay in prose",
            "Left mock cart — browser-local qty/subtotal/thumbs via /api/part-image",
        ],
    )

    add_bullets(
        prs,
        "Agentic architecture",
        [
            "Tool loop with hard round cap — chain normalize → lookup / compat / symptom",
            "Catalog = source of truth — tool output + retrieveExact → retrieval → cards",
            "One LLM path for catalog — no API key → 503 on catalog turns; hello/glossary/OOS without model",
            "Session memory (gated) — SessionContext + system note; allowSessionCarryForRetrieval for retrieval & prompt",
        ],
    )

    add_bullets(
        prs,
        "Extensibility & scalability",
        [
            "catalog.json — diff-friendly demo DB",
            "retrieveExact — single hybrid retrieval API",
            "Optional RAG — semantic_search + embeddings.json + in-memory cosine",
        ],
    )

    add_bullets(
        prs,
        "Six query categories → tools → UI",
        [
            "1 Install — get_install_guide, lookup_part → support·install",
            "2 Compatibility — check_compatibility, normalize → support·compat",
            "3 Symptom/repair — search_by_symptom, catalog_search → repair / candidates",
            "4 Lookup & browse — lookup_part, catalog_search → product / candidates",
            "5 Out-of-scope — no tools → refusal, no block",
            "6 Clarify — optional normalize → question + chips, no block",
            "No block order: OOS → clarify → no-match copy",
        ],
    )

    add_bullets(
        prs,
        "Tool catalog",
        [
            "normalize_part_number · lookup_part · check_compatibility · get_install_guide",
            "search_by_symptom · semantic_search · fetch_part_page · catalog_search",
        ],
    )

    add_bullets(
        prs,
        "Session & scope",
        [
            "System prompt → deterministic OOS gate → no_evidence + alignCitationsToBlocks",
            "Clarify UX: Example: lines → chips; other tips in reply body",
        ],
    )

    add_bullets(
        prs,
        "Server pipeline",
        [
            "POST /api/chat → glossary / conversation shortcuts → OOS gate",
            "→ LLM tool loop → buildBlocksFromRetrieval → clarify if no blocks",
            "→ alignCitationsToBlocks → done { blocks, citations, suggested_actions, tool_trace, … }",
        ],
    )

    add_bullets(
        prs,
        "Data artifacts",
        [
            "web/data/catalog.json — source of truth",
            "web/data/embeddings.json — optional vectors",
            "web/scripts/scrape-parts.mjs · generate-embeddings.mjs",
        ],
    )

    add_bullets(
        prs,
        "Invariants & evaluation",
        [
            "Structured facts from catalog/tools — not hallucinated",
            "No API key: paths per DESIGN.md (glossary / OOS / etc.)",
            "Stream closes in finally",
            "Golden: web/scripts/goldenCases.mjs — tool_trace assertions when LLM on",
        ],
    )

    add_title_slide(prs, "Thank you", "Source: DESIGN.md · Regenerate: python3 docs/generate_design_deck.py")

    out = Path(__file__).resolve().parent / "PartSelect-Design-Deck.pptx"
    prs.save(out)
    print(f"Wrote {out}")


if __name__ == "__main__":
    main()
